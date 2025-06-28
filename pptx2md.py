from pptx import Presentation
from pptx.shapes.autoshape import Shape as AutoShape

prs = Presentation('test.pptx')

with open('output.md', 'w', encoding='utf-8') as f:
    for slide in prs.slides:
        for shape in slide.shapes:
            
            if isinstance(shape, AutoShape) and shape.has_text_frame:
                text = shape.text_frame.text
                if text:
                    f.write(text + '\n\n')  
